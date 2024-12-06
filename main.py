import os
import json
import yaml
import mimetypes
import xml.etree.ElementTree as ET
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation
from PIL import Image
import base64
import csv
import PyPDF2


class UnifiedFileReader:
    MIME_TYPES = {
        ".txt": "text/plain",
        ".json": "application/json",
        ".yaml": "text/x-yaml",
        ".yml": "text/x-yaml",
        ".xml": "application/xml",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".csv": "text/csv",
        ".pdf": "application/pdf",
        ".png": "image/png",
        ".jpeg": "image/jpeg",
        ".jpg": "image/jpeg",
        ".py": "text/x-python",
        ".js": "text/javascript",
        ".java": "text/x-java-source",
        ".cpp": "text/x-c++src",
        ".html": "text/html",
    }

    def read_file(self, file_path):
        """Determine file type and read the file accordingly."""
        file_extension = os.path.splitext(file_path)[1].lower()
        file_type = self.MIME_TYPES.get(file_extension, None)

        if file_type is None:
            raise ValueError(f"Unsupported file type: {file_extension}")

        if file_type in [
            "text/plain",
            "text/x-python",
            "text/javascript",
            "text/x-java-source",
            "text/x-c++src",
            "text/html",
        ]:
            return self.read_text(file_path)
        elif file_type == "application/json":
            return self.read_json(file_path)
        elif file_type in ["text/x-yaml", "application/xml"]:
            return self.read_yaml_or_xml(file_path, file_type)
        elif (
            file_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            return self.read_docx(file_path)
        elif (
            file_type
            == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            return self.read_excel(file_path)
        elif (
            file_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            return self.read_pptx(file_path)
        elif file_type == "text/csv":
            return self.read_csv(file_path)
        elif file_type == "application/pdf":
            return self.read_pdf(file_path)
        elif file_type in ["image/png", "image/jpeg"]:
            return self.read_image(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    @staticmethod
    def read_text(file_path):
        """Read plain text and code files."""
        with open(file_path, "r", encoding="utf-8") as file:
            content = file.read()
        return {"file_name": os.path.basename(file_path), "content": content}

    @staticmethod
    def read_json(file_path):
        """Read JSON files."""
        with open(file_path, "r", encoding="utf-8") as file:
            content = json.load(file)
        return {"file_name": os.path.basename(file_path), "content": content}

    @staticmethod
    def read_yaml_or_xml(file_path, file_type):
        """Read YAML or XML files."""
        with open(file_path, "r", encoding="utf-8") as file:
            content = file.read()
        if file_type == "application/xml":
            try:
                root = ET.fromstring(content)
                content = UnifiedFileReader.xml_to_dict(root)
            except ET.ParseError as e:
                raise ValueError(f"Error parsing XML: {str(e)}")
        elif file_type == "text/x-yaml":
            content = yaml.safe_load(content)
        return {"file_name": os.path.basename(file_path), "content": content}

    @staticmethod
    def xml_to_dict(element):
        """Convert XML element to a dictionary."""
        return {element.tag: {child.tag: child.text for child in element}}

    @staticmethod
    def read_docx(file_path):
        """Read DOCX files."""
        doc = Document(file_path)
        content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return {"file_name": os.path.basename(file_path), "content": content}

    @staticmethod
    def read_excel(file_path):
        """Read Excel files."""
        workbook = load_workbook(file_path)
        sheet = workbook.active
        content = [[cell.value for cell in row] for row in sheet.iter_rows()]
        return {"file_name": os.path.basename(file_path), "content": content}

    @staticmethod
    def read_pptx(file_path):
        """Read PPTX files."""
        presentation = Presentation(file_path)
        content = []
        for slide in presentation.slides:
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_content.append(shape.text)
            content.append(slide_content)
        return {"file_name": os.path.basename(file_path), "content": content}

    @staticmethod
    def read_csv(file_path):
        """Read CSV files."""
        with open(file_path, "r", encoding="utf-8") as file:
            reader = csv.reader(file)
            content = [row for row in reader]
        return {"file_name": os.path.basename(file_path), "content": content}

    @staticmethod
    def read_pdf(file_path):
        """Read PDF files."""
        content = []
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                content.append(page.extract_text())
        return {"file_name": os.path.basename(file_path), "content": "\n".join(content)}

    @staticmethod
    def read_image(file_path):
        """Read image files and return base64-encoded data."""
        with open(file_path, "rb") as img_file:
            encoded_string = base64.b64encode(img_file.read()).decode("utf-8")
        return {"file_name": os.path.basename(file_path), "content": encoded_string}

    @staticmethod
    def write_output(data, output_file="output.json"):
        """Write the output data to a JSON file."""
        try:
            formatted_content = data["content"]
            if isinstance(formatted_content, str):
                data["content"] = formatted_content.splitlines()

            with open(output_file, "w", encoding="utf-8") as outfile:
                json.dump(data, outfile, indent=4, ensure_ascii=False)

            print(f"Data written to {output_file}")
        except Exception as e:
            raise ValueError(f"Error writing output file: {str(e)}")


# Example Usage
if __name__ == "__main__":
    file_name = "ch3.pdf"
    reader = UnifiedFileReader()

    # Read the file
    try:
        data = reader.read_file(file_name)

        # Write the output to a JSON file
        reader.write_output(data)
    except ValueError as e:
        print(f"Error: {str(e)}")